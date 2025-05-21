from flask import Flask, request, jsonify, send_file
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image
import requests
from transformers import BlipProcessor, BlipForConditionalGeneration
from TTS.api import TTS
import os

app = Flask(__name__)
processor = BlipProcessor.from_pretrained("Salesforce/blip-image-captioning-base")
model = BlipForConditionalGeneration.from_pretrained("Salesforce/blip-image-captioning-base")
tts = TTS(model_name="tts_models/en/ljspeech/vits", progress_bar=False, gpu=False)

UPLOAD_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'uploads')
ALLOWED_EXTENSIONS = {'ppt', 'pptx'}
UPLOADED_IMAGES_FOLDER = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'output', 'images')


uploaded_file_name = None

if not os.path.exists(UPLOAD_FOLDER):
    os.makedirs(UPLOAD_FOLDER)

@app.after_request
def after_request(response):
    response.headers.add('Access-Control-Allow-Origin', '*')
    response.headers.add('Access-Control-Allow-Headers', 'Content-Type,Authorization')
    response.headers.add('Access-Control-Allow-Methods', 'GET,PUT,POST,DELETE,OPTIONS')
    return response

# @app.route('/test', methods=['GET', 'OPTIONS'])
# def test():
#     if request.method == 'OPTIONS':
#         return '', 200
#     return jsonify({'message': 'API is working!'}), 200

@app.route('/upload', methods=['POST', 'OPTIONS'])
def upload_file():
    if request.method == 'OPTIONS':
        return '', 200
        
    print("Upload request received")
    
    if 'file' not in request.files:
        return jsonify({'message': 'No file part'}), 400
        
    file = request.files['file']
    
    if file.filename == '':
        return jsonify({'message': 'No selected file'}), 400
        
    if file and '.' in file.filename and file.filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS:
        filename = os.path.join(UPLOAD_FOLDER, file.filename)
        file.save(filename)
        extract_ppt_content(filename)    
        return jsonify({'message': 'File successfully uploaded', 'filename': filename}), 200
    else:
        return jsonify({'message': 'File type not allowed'}), 400
    
# def convert_ppt_to_pptx(file_path):
#     powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
#     powerpoint.Visible = 1
    
#     ppt_path = os.path.abspath(file_path)
#     pptx_path = ppt_path + "x"
    
#     presentation = powerpoint.Presentations.Open(ppt_path)
#     presentation.SaveAs(pptx_path, FileFormat=24)
#     presentation.Close()
#     powerpoint.Quit()
#     return pptx_path
    
def extract_ppt_content(file_path, output_dir="output", output_text_file="output.txt"):
    presentation = Presentation(file_path)
    os.makedirs(output_dir, exist_ok=True)
    
    images_dir = os.path.join(output_dir, "images")
    os.makedirs(images_dir, exist_ok=True)
    
    lines = []
    
    for slide_index, slide in enumerate(presentation.slides, start=1):
        lines.append(f"\n--- Slide {slide_index} ---\n")
        for shape_index, shape in enumerate(slide.shapes):
            if hasattr(shape, "text"):
                lines.append(shape.text.strip())
                
            if(shape.has_table):
                for row in shape.table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    lines.append(" | ".join(row_text))
                    
            if(shape.shape_type == MSO_SHAPE_TYPE.PICTURE):
                image = shape.image
                image_bytes = image.blob
                extension = image.ext
                image_filename = f"slide({slide_index}_img{shape_index}).{extension}"
                image_path = os.path.join(images_dir, image_filename)
                
                with open(image_path, "wb") as file:
                    file.write(image_bytes)
                
                lines.append(f"[Image: {image_filename}]")
                
    output_text_file_path = os.path.join(output_dir, output_text_file)
    with open(output_text_file_path, "w", encoding="utf-8") as file:
        file.write("\n".join(lines))
        
    global uploaded_file_name
    uploaded_file_name = output_text_file_path
    
    print(f"Text and image references saved to {output_text_file_path}")
    print(f"Images saved in {images_dir}")
    
def generate_image_caption(image_path: str) -> str:
    image = Image.open(image_path).convert("RGB")
    inputs = processor(images=image, return_tensors="pt")
    out = model.generate(**inputs)
    caption = processor.decode(out[0], skip_special_tokens=True)
    return caption


def process_images_in_folder(folder_path: str):
    captions = []
    
    for index, filename in enumerate(os.listdir(folder_path)):
        if filename.endswith(('.png', '.jpg', '.jpeg')):
            image_path = os.path.join(folder_path, filename)
            image_caption = generate_image_caption(image_path)
            captions.append(f"[Image {index} caption: {image_caption}]")
            
    return captions
            
@app.route('/generate', methods=['GET'])
def generate_podcast_script():
    input_image_caption = process_images_in_folder(UPLOADED_IMAGES_FOLDER)
    with open(uploaded_file_name, 'r') as file:
        input_text = file.read()

        
    url = "http://localhost:11434/api/generate"
    prompt = (
        " There should be no headers or subtitles, just the text to be spoken in a single paragraph. Remove all Markdown formatting from the following text, including asterisks used for bold or italic, and hashtags used for headings. Keep the text content and structure intact. Based on the following article and image description, generate an educational podcast script with a natural tone. The podcast refelcts a teacher teaching. The script should be insightful and include examples to grasp the material.\n\n"
        f"Article:\n{input_text}\n\n"
        f"Image Description:\n{input_image_caption}"
    )
    
    payload = {
        "model": "gemma",
        "prompt": prompt,
        "stream": False
    }
    
    response = requests.post(url, json=payload)
    print(f"Response: {response.json().get('response')}")
    audio_path = os.path.join(os.path.dirname(__file__), "podcast_output.wav")
    generate_podcast_audio(response.json().get("response"))
    return send_file(audio_path, mimetype='audio/wav', as_attachment=True)

def generate_podcast_audio(input_text: str, output_path="podcast_output.wav"):
    tts.tts_to_file(text=input_text, file_path=output_path)
    


if __name__ == '__main__':
    print(f"Starting server... Upload folder: {UPLOAD_FOLDER}")
    PORT = 8000
    print(f"Server running on http://localhost:{PORT}")
    app.run(debug=True, host='0.0.0.0', port=PORT)